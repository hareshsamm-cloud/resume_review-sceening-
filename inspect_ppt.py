from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ppt_path = r"C:\Users\Hares\Downloads\Software Hackathon PPT Template.pptx"

def inspect_shape(shape, prefix=""):
    shape_type = shape.shape_type
    name = shape.name
    text = ""
    if shape.has_text_frame:
        text = shape.text.strip().replace('\n', ' | ')
    
    print(f"{prefix}Shape: Name='{name}', Type={shape_type}, Text='{text[:100]}'")
    
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for idx, sub_shape in enumerate(shape.shapes):
            inspect_shape(sub_shape, prefix + f"  [{idx}] ")

try:
    prs = Presentation(ppt_path)
    print(f"Total Slides: {len(prs.slides)}")
    
    for idx, slide in enumerate(prs.slides):
        print(f"\n--- Slide {idx + 1} ---")
        for s_idx, shape in enumerate(slide.shapes):
            inspect_shape(shape, f"S{s_idx} ")
            
except Exception as e:
    print(f"Error reading presentation: {e}")
