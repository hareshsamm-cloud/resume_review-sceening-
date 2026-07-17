from pptx import Presentation
from pptx.util import Inches

ppt_path = r"C:\Users\Hares\Downloads\Software Hackathon PPT Template.pptx"
prs = Presentation(ppt_path)
slide = prs.slides[0]

for shape in slide.shapes:
    if shape.has_text_frame:
        left = shape.left / Inches(1)
        top = shape.top / Inches(1)
        width = shape.width / Inches(1)
        height = shape.height / Inches(1)
        print(f"Name: {shape.name}")
        print(f"  Position: Left={left:.2f}\", Top={top:.2f}\", Width={width:.2f}\", Height={height:.2f}\"")
        print(f"  Text: {shape.text}")
