import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

ppt_path = r"C:\Users\Hares\Downloads\Software Hackathon PPT Template.pptx"

def set_shape_transparency(shape, transparency_fraction):
    """Sets shape transparency (0.0 = opaque, 1.0 = fully transparent) via raw XML."""
    try:
        spPr = shape._element.spPr
        solidFill = spPr.xpath('a:solidFill')
        if solidFill:
            color_elem = solidFill[0]
            color_clrs = color_elem.xpath('a:srgbClr') or color_elem.xpath('a:schemeClr')
            if color_clrs:
                clr = color_clrs[0]
                # Clear existing alpha tags
                for alpha in clr.xpath('a:alpha'):
                    clr.remove(alpha)
                # Append alpha element
                opacity_val = int((1.0 - transparency_fraction) * 100000)
                alpha = OxmlElement('a:alpha')
                alpha.set('val', str(opacity_val))
                clr.append(alpha)
    except Exception as e:
        print(f"Error setting transparency: {e}")

try:
    prs = Presentation(ppt_path)
    width_in = prs.slide_width / Inches(1)
    height_in = prs.slide_height / Inches(1)
    print(f"Slide Dimensions: Width={width_in:.2f}\", Height={height_in:.2f}\"")
    
    # -------------------------------------------------------------
    # SLIDE 1: Problem and Objectives (Large fonts, wide spread)
    # -------------------------------------------------------------
    slide1 = prs.slides[0]
    
    # 1. Clean up any existing Title box we added on Slide 1
    for shape in list(slide1.shapes):
        if shape.name == "TopicTitleBox":
            sp = shape._element
            sp.getparent().remove(sp)
            
    # 2. Add Topic Title box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(18.0), Inches(2.2))
    title_box.name = "TopicTitleBox"
    tf = title_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "Topic: ResumeSphere AI"
    p1.font.bold = True
    p1.font.size = Pt(48)
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.font.name = "Outfit"
    p1.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = "Next-Gen Multi-Portal Resume Screener, Integrity Auditor & Career Path Recommendation System"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(6, 182, 212) # Cyan accent
    p2.font.name = "Outfit"
    
    tb13 = None
    tb14 = None
    for shape in slide1.shapes:
        if shape.name == 'TextBox 13':
            tb13 = shape
        elif shape.name == 'TextBox 14':
            tb14 = shape
            
    if tb13:
        tb13.left = Inches(1.0)
        tb13.top = Inches(6.2)
        tb13.width = Inches(8.5)
        tb13.height = Inches(4.5)
        
        tf = tb13.text_frame
        tf.clear()
        tf.word_wrap = True
        
        # Add Team Name
        p1 = tf.paragraphs[0]
        p1.text = "Team Name: "
        run1 = p1.add_run()
        run1.text = "Data Breakers"
        run1.font.bold = True
        run1.font.color.rgb = RGBColor(59, 130, 246) # Blue accent
        p1.font.size = Pt(24)
        p1.space_after = Pt(8)
        
        # Add Team Leader
        p2 = tf.add_paragraph()
        p2.text = "Team Leader: "
        run2 = p2.add_run()
        run2.text = "Abinav Balasubramaniam"
        run2.font.bold = True
        p2.font.size = Pt(24)
        p2.space_after = Pt(20)
        
        # Add Problem Statement
        p3 = tf.add_paragraph()
        p3.text = "Problem Statement:"
        p3.font.bold = True
        p3.font.size = Pt(22)
        p3.font.color.rgb = RGBColor(255, 255, 255)
        p3.space_after = Pt(8)
        
        p4 = tf.add_paragraph()
        p4.text = "Evaluating candidate resumes manually is highly time-consuming, prone to human cognitive bias, and offers no career up-skilling guides to help students identify and resolve missing required skills."
        p4.font.size = Pt(17)
        p4.font.color.rgb = RGBColor(210, 215, 230)
        
    if tb14:
        # Position tb14 side-by-side on the right half of Slide 1
        tb14.left = Inches(10.5)
        tb14.top = Inches(6.2)
        tb14.width = Inches(8.5)
        tb14.height = Inches(4.5)
        
        tf14 = tb14.text_frame
        tf14.clear()
        tf14.word_wrap = True
        
        p_obj_head = tf14.paragraphs[0]
        p_obj_head.text = "Project Objectives:"
        p_obj_head.font.bold = True
        p_obj_head.font.size = Pt(24)
        p_obj_head.font.color.rgb = RGBColor(139, 92, 246) # Purple accent
        p_obj_head.space_after = Pt(16)
        
        objectives = [
            "Implement an automated, 100% offline, resume parser and scorer using Python and PyPDF.",
            "Provide dual recruiter and student dashboards (Leaderboard vs. Personal Bio Highlights).",
            "Design a career recommendation system matching top companies and providing skills gap roadmaps."
        ]
        
        for obj in objectives:
            p_obj = tf14.add_paragraph()
            p_obj.text = f"• {obj}"
            p_obj.font.size = Pt(17)
            p_obj.space_after = Pt(12)
            p_obj.font.color.rgb = RGBColor(210, 215, 230)

    # -------------------------------------------------------------
    # SLIDE 2: 20-Hour Roadmap (Cleanup overlaps & set XML transparency)
    # -------------------------------------------------------------
    slide2 = prs.slides[1]
    
    # 1. Clean up ALL shapes we added in past runs
    background_shapes = {"Freeform 2", "Freeform 3", "Group 4"}
    shapes_to_remove = []
    for shape in slide2.shapes:
        if shape.name not in background_shapes:
            shapes_to_remove.append(shape)
            
    print(f"Removing {len(shapes_to_remove)} added shapes from Slide 2...")
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
        
    # 2. Add title text box at the top
    title_box = slide2.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(12.0), Inches(0.8))
    title_box.name = "RoadmapTitle"
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "20-Hour Development Roadmap"
    p_title.font.bold = True
    p_title.font.size = Pt(36)
    p_title.font.color.rgb = RGBColor(255, 255, 255)
    
    # Roadmap Phase Cards
    phases = [
        {
            "hours": "Hours 0 - 4",
            "title": "Phase 1: Setup & Design",
            "points": [
                "Initialize Git and .gitignore.",
                "Django project & app structures.",
                "Database schema configuration.",
                "Visual CSS design variables."
            ],
            "color": RGBColor(59, 130, 246) # Blue
        },
        {
            "hours": "Hours 4 - 10",
            "title": "Phase 2: Parser & Logic",
            "points": [
                "Implement PDF text extraction.",
                "Write parser regex engines.",
                "Set up tech skills map.",
                "Integrate eligibility rules."
            ],
            "color": RGBColor(6, 182, 212) # Cyan
        },
        {
            "hours": "Hours 10 - 16",
            "title": "Phase 3: Interfaces",
            "points": [
                "Recruiter upload & cockpit.",
                "Student dashboard layouts.",
                "Integrate 50+ job role menus.",
                "Client-side list search filters."
            ],
            "color": RGBColor(139, 92, 246) # Purple
        },
        {
            "hours": "Hours 16 - 20",
            "title": "Phase 4: Integration & QA",
            "points": [
                "Mock email outgoing log outbox.",
                "Verify strict code compilings.",
                "Sync conflicts & push to Git.",
                "System flow check & QA."
            ],
            "color": RGBColor(16, 185, 129) # Green
        }
    ]
    
    col_width = Inches(4.3)
    gap = Inches(0.5)
    top_pos = Inches(3.2)
    
    for idx, phase in enumerate(phases):
        left_pos = Inches(1.0) + idx * (col_width + gap)
        
        # Add shape cards (a rectangle to act as a backing)
        card = slide2.shapes.add_shape(
            1, # Rectangle (MSO_SHAPE.RECTANGLE)
            left_pos, top_pos, col_width, Inches(6.0)
        )
        card.name = f"RoadmapCard_{idx}"
        
        # Style the shape background to be transparent
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(17, 21, 44)
        
        # Apply XML-level transparency (40% transparent / 60% opacity)
        set_shape_transparency(card, 0.40)
            
        card.line.color.rgb = phase["color"]
        card.line.width = Pt(2.5)
        
        # Write text inside the card shape
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.25)
        tf_card.margin_right = Inches(0.25)
        tf_card.margin_top = Inches(0.35)
        tf_card.margin_bottom = Inches(0.3)
        
        # Hours Header
        p_hrs = tf_card.paragraphs[0]
        p_hrs.text = phase["hours"]
        p_hrs.font.bold = True
        p_hrs.font.size = Pt(20)
        p_hrs.font.color.rgb = phase["color"]
        p_hrs.space_after = Pt(4)
        
        # Phase Title
        p_p_title = tf_card.add_paragraph()
        p_p_title.text = phase["title"]
        p_p_title.font.bold = True
        p_p_title.font.size = Pt(17)
        p_p_title.space_after = Pt(14)
        p_p_title.font.color.rgb = RGBColor(255, 255, 255)
        
        # Phase bullet points
        for pt in phase["points"]:
            p_pt = tf_card.add_paragraph()
            p_pt.text = f"• {pt}"
            p_pt.font.size = Pt(13.5)
            p_pt.space_after = Pt(10)
            p_pt.font.color.rgb = RGBColor(200, 210, 225)
            
    # -------------------------------------------------------------
    # SLIDE 3: Improvements Made After Review 1
    # -------------------------------------------------------------
    slide3 = prs.slides[2]
    tb3_title = None
    ic0 = None
    ic1 = None
    ic2 = None
    for shape in slide3.shapes:
        if shape.name == 'TextBox 3':
            tb3_title = shape
        elif shape.name == 'ImprovementCard_0':
            ic0 = shape
        elif shape.name == 'ImprovementCard_1':
            ic1 = shape
        elif shape.name == 'ImprovementCard_2':
            ic2 = shape
            
    if tb3_title:
        tb3_title.text_frame.clear()
        p = tb3_title.text_frame.paragraphs[0]
        p.text = "Improvements Made After Review 1"
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = "Outfit"
        
    # Re-create and style Improvement Card 0 (removes template theme overrides)
    if ic0:
        left, top, width, height = ic0.left, ic0.top, ic0.width, ic0.height
        sp = ic0._element
        sp.getparent().remove(sp)
        
        ic0 = slide3.shapes.add_shape(1, left, top, width, height)
        ic0.name = 'ImprovementCard_0'
        ic0.fill.solid()
        ic0.fill.fore_color.rgb = RGBColor(17, 21, 44)
        set_shape_transparency(ic0, 0.40)
        ic0.line.color.rgb = RGBColor(59, 130, 246) # Blue
        ic0.line.width = Pt(2.5)
        
        tf0 = ic0.text_frame
        tf0.clear()
        tf0.word_wrap = True
        tf0.margin_left = Inches(0.25)
        tf0.margin_right = Inches(0.25)
        tf0.margin_top = Inches(0.35)
        
        p_hdr = tf0.paragraphs[0]
        p_hdr.text = "01 | Dynamic Recalculation Engine"
        p_hdr.font.bold = True
        p_hdr.font.size = Pt(19)
        p_hdr.font.color.rgb = RGBColor(59, 130, 246)
        p_hdr.space_after = Pt(14)
        
        points = [
            "Replaced static matching with real-time scoring audits.",
            "Toggling experience sliders or templates recalculates all matches instantly.",
            "Fully preserves state using Django session cache parameters."
        ]
        for pt in points:
            p_pt = tf0.add_paragraph()
            p_pt.text = f"• {pt}"
            p_pt.font.size = Pt(13.5)
            p_pt.space_after = Pt(10)
            p_pt.font.color.rgb = RGBColor(200, 210, 225)
            p_pt.font.name = "Outfit"
            
    # Re-create and style Improvement Card 1
    if ic1:
        left, top, width, height = ic1.left, ic1.top, ic1.width, ic1.height
        sp = ic1._element
        sp.getparent().remove(sp)
        
        ic1 = slide3.shapes.add_shape(1, left, top, width, height)
        ic1.name = 'ImprovementCard_1'
        ic1.fill.solid()
        ic1.fill.fore_color.rgb = RGBColor(17, 21, 44)
        set_shape_transparency(ic1, 0.40)
        ic1.line.color.rgb = RGBColor(139, 92, 246) # Purple
        ic1.line.width = Pt(2.5)
        
        tf1 = ic1.text_frame
        tf1.clear()
        tf1.word_wrap = True
        tf1.margin_left = Inches(0.25)
        tf1.margin_right = Inches(0.25)
        tf1.margin_top = Inches(0.35)
        
        p_hdr = tf1.paragraphs[0]
        p_hdr.text = "02 | 16-Tier Anomaly & Fraud Audit"
        p_hdr.font.bold = True
        p_hdr.font.size = Pt(19)
        p_hdr.font.color.rgb = RGBColor(139, 92, 246)
        p_hdr.space_after = Pt(14)
        
        points = [
            "Detects keyword stuffing (>20 skills) and duplicate word loops.",
            "Flags certification timeline release gaps (e.g. React in 2008).",
            "Identifies multi-stack tech domain conflicts and blank metadata."
        ]
        for pt in points:
            p_pt = tf1.add_paragraph()
            p_pt.text = f"• {pt}"
            p_pt.font.size = Pt(13.5)
            p_pt.space_after = Pt(10)
            p_pt.font.color.rgb = RGBColor(200, 210, 225)
            p_pt.font.name = "Outfit"

    # Re-create and style Improvement Card 2
    if ic2:
        left, top, width, height = ic2.left, ic2.top, ic2.width, ic2.height
        sp = ic2._element
        sp.getparent().remove(sp)
        
        ic2 = slide3.shapes.add_shape(1, left, top, width, height)
        ic2.name = 'ImprovementCard_2'
        ic2.fill.solid()
        ic2.fill.fore_color.rgb = RGBColor(17, 21, 44)
        set_shape_transparency(ic2, 0.40)
        ic2.line.color.rgb = RGBColor(16, 185, 129) # Green
        ic2.line.width = Pt(2.5)
        
        tf2 = ic2.text_frame
        tf2.clear()
        tf2.word_wrap = True
        tf2.margin_left = Inches(0.25)
        tf2.margin_right = Inches(0.25)
        tf2.margin_top = Inches(0.35)
        
        p_hdr = tf2.paragraphs[0]
        p_hdr.text = "03 | Outbox SMTP Simulator & Space"
        p_hdr.font.bold = True
        p_hdr.font.size = Pt(19)
        p_hdr.font.color.rgb = RGBColor(16, 185, 129)
        p_hdr.space_after = Pt(14)
        
        points = [
            "Visual decision log captures sent offer and rejection emails.",
            "Created offline College Placement deck with 210 companies.",
            "Responsive searching and filtering mapped by Green, Orange, and Red tiers."
        ]
        for pt in points:
            p_pt = tf2.add_paragraph()
            p_pt.text = f"• {pt}"
            p_pt.font.size = Pt(13.5)
            p_pt.space_after = Pt(10)
            p_pt.font.color.rgb = RGBColor(200, 210, 225)
            p_pt.font.name = "Outfit"

    # -------------------------------------------------------------
    # SLIDE 4: Core Implemented Systems
    # -------------------------------------------------------------
    slide4 = prs.slides[3]
    tb3 = None
    card0 = None
    card1 = None
    card2 = None
    for shape in slide4.shapes:
        if shape.name == 'TextBox 3':
            tb3 = shape
        elif shape.name == 'JourneyCard_0':
            card0 = shape
        elif shape.name == 'JourneyCard_1':
            card1 = shape
        elif shape.name == 'JourneyCard_2':
            card2 = shape
            
    if tb3:
        tb3.text_frame.clear()
        p = tb3.text_frame.paragraphs[0]
        p.text = "ResumeSphere AI - Core Implemented Systems"
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = "Outfit"
        
    # Re-create and style Card 0 (removes template theme overrides)
    if card0:
        left, top, width, height = card0.left, card0.top, card0.width, card0.height
        sp = card0._element
        sp.getparent().remove(sp)
        
        card0 = slide4.shapes.add_shape(1, left, top, width, height)
        card0.name = 'JourneyCard_0'
        card0.fill.solid()
        card0.fill.fore_color.rgb = RGBColor(17, 21, 44)
        set_shape_transparency(card0, 0.40)
        card0.line.color.rgb = RGBColor(59, 130, 246) # Blue
        card0.line.width = Pt(2.5)
        
        tf0 = card0.text_frame
        tf0.clear()
        tf0.word_wrap = True
        tf0.margin_left = Inches(0.25)
        tf0.margin_right = Inches(0.25)
        tf0.margin_top = Inches(0.35)
        
        p_hdr = tf0.paragraphs[0]
        p_hdr.text = "01 | Multi-Portal Workspace"
        p_hdr.font.bold = True
        p_hdr.font.size = Pt(19)
        p_hdr.font.color.rgb = RGBColor(59, 130, 246)
        p_hdr.space_after = Pt(14)
        
        points = [
            "Recruiter Cockpit: Visual leaderboards, bulk actions, and direct candidate processing.",
            "Student Space: Resume upload analytics, company eligibility status, and career guidelines.",
            "College Admin Portal: Configure blacklist limits, check fraud logs, and manage suspensions."
        ]
        for pt in points:
            p_pt = tf0.add_paragraph()
            p_pt.text = f"• {pt}"
            p_pt.font.size = Pt(13)
            p_pt.space_after = Pt(10)
            p_pt.font.color.rgb = RGBColor(200, 210, 225)
            p_pt.font.name = "Outfit"
            
    # Re-create and style Card 1
    if card1:
        left, top, width, height = card1.left, card1.top, card1.width, card1.height
        sp = card1._element
        sp.getparent().remove(sp)
        
        card1 = slide4.shapes.add_shape(1, left, top, width, height)
        card1.name = 'JourneyCard_1'
        card1.fill.solid()
        card1.fill.fore_color.rgb = RGBColor(17, 21, 44)
        set_shape_transparency(card1, 0.40)
        card1.line.color.rgb = RGBColor(139, 92, 246) # Purple
        card1.line.width = Pt(2.5)
        
        tf1 = card1.text_frame
        tf1.clear()
        tf1.word_wrap = True
        tf1.margin_left = Inches(0.25)
        tf1.margin_right = Inches(0.25)
        tf1.margin_top = Inches(0.35)
        
        p_hdr = tf1.paragraphs[0]
        p_hdr.text = "02 | Digital Resume Fingerprint"
        p_hdr.font.bold = True
        p_hdr.font.size = Pt(19)
        p_hdr.font.color.rgb = RGBColor(139, 92, 246)
        p_hdr.space_after = Pt(14)
        
        points = [
            "6-Tier Anomaly Audit: Score Quality, Authenticity, Fraud, Communication, Tech Depth, and Learning.",
            "Integrity Auditing: Hidden font parsing, background keyword stuffing detection, and email domain check.",
            "Anti-Gaming Blocker: Automated college reports and warning thresholds preventing student profile abuse."
        ]
        for pt in points:
            p_pt = tf1.add_paragraph()
            p_pt.text = f"• {pt}"
            p_pt.font.size = Pt(13)
            p_pt.space_after = Pt(10)
            p_pt.font.color.rgb = RGBColor(200, 210, 225)
            p_pt.font.name = "Outfit"

    # Re-create and style Card 2
    if card2:
        left, top, width, height = card2.left, card2.top, card2.width, card2.height
        sp = card2._element
        sp.getparent().remove(sp)
        
        card2 = slide4.shapes.add_shape(1, left, top, width, height)
        card2.name = 'JourneyCard_2'
        card2.fill.solid()
        card2.fill.fore_color.rgb = RGBColor(17, 21, 44)
        set_shape_transparency(card2, 0.40)
        card2.line.color.rgb = RGBColor(16, 185, 129) # Green
        card2.line.width = Pt(2.5)
        
        tf2 = card2.text_frame
        tf2.clear()
        tf2.word_wrap = True
        tf2.margin_left = Inches(0.25)
        tf2.margin_right = Inches(0.25)
        tf2.margin_top = Inches(0.35)
        
        p_hdr = tf2.paragraphs[0]
        p_hdr.text = "03 | Stress Test & Battles"
        p_hdr.font.bold = True
        p_hdr.font.size = Pt(19)
        p_hdr.font.color.rgb = RGBColor(16, 185, 129)
        p_hdr.space_after = Pt(14)
        
        points = [
            "Explainability Ledger: Visual green (+)/red (-) ledger detailing points logic behind every match.",
            "ATS Stress Testing: Programmatic evaluation maps against Google, Amazon, Stripe, Microsoft, Zoho.",
            "Candidate Battles: Dual-column head-to-head matching overlay declaring crowned winner."
        ]
        for pt in points:
            p_pt = tf2.add_paragraph()
            p_pt.text = f"• {pt}"
            p_pt.font.size = Pt(13)
            p_pt.space_after = Pt(10)
            p_pt.font.color.rgb = RGBColor(200, 210, 225)
            p_pt.font.name = "Outfit"

    # -------------------------------------------------------------
    # SLIDE 6: Thank You details
    # -------------------------------------------------------------
    slide6 = prs.slides[5]
    for shape in list(slide6.shapes):
        if shape.name == "ThankYouSubtitle":
            sp = shape._element
            sp.getparent().remove(sp)
            
    sub_box = slide6.shapes.add_textbox(Inches(1.0), Inches(8.2), Inches(18.0), Inches(1.5))
    sub_box.name = "ThankYouSubtitle"
    tf6 = sub_box.text_frame
    tf6.word_wrap = True
    p6 = tf6.paragraphs[0]
    p6.text = "Questions & Discussion"
    p6.alignment = 1 # Center alignment
    p6.font.bold = True
    p6.font.size = Pt(24)
    p6.font.color.rgb = RGBColor(6, 182, 212) # Cyan accent
    p6.font.name = "Outfit"
    p6.space_after = Pt(8)
    
    p6_team = tf6.add_paragraph()
    p6_team.text = "Data Breakers Team  |  Abinav Balasubramaniam (Leader)"
    p6_team.alignment = 1 # Center
    p6_team.font.size = Pt(18)
    p6_team.font.color.rgb = RGBColor(200, 210, 225)
    p6_team.font.name = "Outfit"

    # Save the modified presentation
    prs.save(ppt_path)
    print("\nPowerPoint modified and saved successfully!")
    
except Exception as e:
    print(f"Error modifying presentation: {e}")
